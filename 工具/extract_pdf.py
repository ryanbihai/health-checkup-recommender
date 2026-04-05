#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF和PPTX文档提取脚本
从PDF和PPTX文件中提取文本内容
"""

import os
import json
import re
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Any

# 路径配置
BASE_DIR = Path(r'c:\IT\02 代理人营销工具\agent-customer-management')
INPUT_FILE = BASE_DIR / '客户档案整理项目' / '中间数据' / '文档清单.json'
OUTPUT_DIR = BASE_DIR / '客户档案整理项目' / '中间数据' / '提取数据'

def extract_pdf_text(filepath: Path) -> Dict[str, Any]:
    """从PDF提取文本"""
    result = {
        'filepath': str(filepath.relative_to(BASE_DIR)),
        'filename': filepath.name,
        'file_type': 'pdf',
        'extract_time': datetime.now().isoformat(),
        'success': False,
        'text': '',
        'page_count': 0,
        'summary': ''
    }
    
    try:
        # 尝试使用PyMuPDF
        import fitz  # PyMuPDF
        
        doc = fitz.open(str(filepath))
        result['page_count'] = len(doc)
        
        full_text = []
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                full_text.append(f"=== 第{page_num}页 ===\n{text}")
        
        result['text'] = '\n\n'.join(full_text)
        result['success'] = True
        
        # 生成摘要
        if result['text']:
            result['summary'] = generate_summary(result['text'], filepath.name)
        
        doc.close()
        
    except ImportError:
        # 尝试使用pdfplumber
        try:
            import pdfplumber
            
            with pdfplumber.open(str(filepath)) as pdf:
                result['page_count'] = len(pdf.pages)
                
                full_text = []
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        full_text.append(f"=== 第{page_num}页 ===\n{text}")
                
                result['text'] = '\n\n'.join(full_text)
                result['success'] = True
                
                if result['text']:
                    result['summary'] = generate_summary(result['text'], filepath.name)
                    
        except ImportError:
            result['error'] = '未安装PDF处理库 (pip install PyMuPDF 或 pip install pdfplumber)'
        except Exception as e:
            result['error'] = str(e)
    except Exception as e:
        result['error'] = str(e)
    
    return result

def extract_pptx_text(filepath: Path) -> Dict[str, Any]:
    """从PPTX提取文本"""
    result = {
        'filepath': str(filepath.relative_to(BASE_DIR)),
        'filename': filepath.name,
        'file_type': 'pptx',
        'extract_time': datetime.now().isoformat(),
        'success': False,
        'text': '',
        'slide_count': 0,
        'slides': [],
        'summary': ''
    }
    
    try:
        from pptx import Presentation
        
        prs = Presentation(str(filepath))
        result['slide_count'] = len(prs.slides)
        
        full_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                slide_content = '\n'.join(slide_text)
                full_text.append(f"=== 第{slide_num}页 ===\n{slide_content}")
                result['slides'].append({
                    'slide_num': slide_num,
                    'text': slide_content
                })
        
        result['text'] = '\n\n'.join(full_text)
        result['success'] = True
        
        if result['text']:
            result['summary'] = generate_summary(result['text'], filepath.name)
        
    except ImportError:
        result['error'] = '未安装PPTX处理库 (pip install python-pptx)'
    except Exception as e:
        result['error'] = str(e)
    
    return result

def generate_summary(text: str, filename: str) -> str:
    """生成文档摘要"""
    # 提取关键信息
    summary = {
        'doc_type': guess_doc_type(text, filename),
        'customer_name': extract_customer_name(text),
        'products': extract_products(text),
        'key_numbers': extract_key_numbers(text),
        'content_preview': text[:500] if text else ''
    }
    
    # 生成摘要文本
    summary_text = []
    
    if summary['doc_type']:
        summary_text.append(f"文档类型: {summary['doc_type']}")
    
    if summary['customer_name']:
        summary_text.append(f"客户姓名: {summary['customer_name']}")
    
    if summary['products']:
        summary_text.append(f"涉及产品: {', '.join(summary['products'][:5])}")
    
    if summary['key_numbers']:
        summary_text.append(f"关键数据: {', '.join(summary['key_numbers'][:3])}")
    
    return '\n'.join(summary_text) if summary_text else '无法生成摘要'

def guess_doc_type(text: str, filename: str) -> str:
    """猜测文档类型"""
    filename_lower = filename.lower()
    text_lower = text.lower()
    
    if '保单' in text or '保单利益' in filename_lower:
        return '保单资料'
    elif '保障方案' in text or '保障规划' in filename_lower:
        return '保障方案'
    elif '产品解析' in text or '产品说明' in filename_lower:
        return '产品解析'
    elif '比较' in text or '对比' in filename_lower:
        return '方案比较'
    elif '电子保单' in filename_lower:
        return '电子保单'
    elif '条款' in text:
        return '保险条款'
    elif '建议书' in text or '建议书' in filename_lower:
        return '保险建议书'
    elif '计划书' in text or '计划书' in filename_lower:
        return '计划书'
    else:
        return '其他文档'

def extract_customer_name(text: str) -> str:
    """提取客户姓名"""
    # 常见模式
    patterns = [
        r'客户[：:]\s*([^\n，。,]+)',
        r'投保人[：:]\s*([^\n，。,]+)',
        r'被保险人[：:]\s*([^\n，。,]+)',
        r'姓名[：:]\s*([^\n，。,]+)',
        r'姓\s*名[：:]\s*([^\n，。,]+)',
    ]
    
    for pattern in patterns:
        match = re.search(pattern, text)
        if match:
            name = match.group(1).strip()
            if len(name) >= 2 and len(name) <= 10:
                return name
    
    return ''

def extract_products(text: str) -> List[str]:
    """提取产品名称"""
    products = set()
    
    # 常见产品关键词
    product_keywords = [
        '大都会', '都会', '康逸', '康健', '安心', '福祥',
        '智行', '尊享', '福寿', '颐年', '赢家', '稳盈'
    ]
    
    for keyword in product_keywords:
        if keyword in text:
            # 提取包含关键词的完整产品名
            pattern = rf'{keyword}[^\n，。,，\s]{{0,10}}'
            matches = re.findall(pattern, text)
            for match in matches:
                if len(match) >= 2:
                    products.add(match)
    
    return list(products)

def extract_key_numbers(text: str) -> List[str]:
    """提取关键数字"""
    numbers = []
    
    # 保费模式
    premium_patterns = [
        r'年缴保费[：:]*\s*([0-9,]+\.?\d*[万千]?元?)',
        r'保费[：:]*\s*([0-9,]+\.?\d*[万千]?元?)',
        r'年缴[：:]*\s*([0-9,]+\.?\d*[万千]?元?)',
    ]
    
    for pattern in premium_patterns:
        matches = re.findall(pattern, text)
        for match in matches:
            if match and len(match) < 20:
                numbers.append(match)
    
    # 保额模式
    amount_patterns = [
        r'保额[：:]*\s*([0-9,]+\.?\d*[万千]?元?)',
        r'保障[：:]*\s*([0-9,]+\.?\d*[万千]?元?)',
    ]
    
    for pattern in amount_patterns:
        matches = re.findall(pattern, text)
        for match in matches:
            if match and len(match) < 20:
                numbers.append(match)
    
    return list(set(numbers))[:10]  # 去重，最多10个

def main():
    """主函数"""
    print("=" * 60)
    print("PDF/PPTX文档提取工具")
    print("=" * 60)
    
    # 确保输出目录存在
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    # 读取文档清单
    print(f"\n📖 读取文档清单...")
    if not INPUT_FILE.exists():
        print(f"❌ 文档清单不存在，请先运行 scan_documents.py")
        return
    
    with open(INPUT_FILE, 'r', encoding='utf-8') as f:
        document_list = json.load(f)
    
    # 筛选PDF和PPTX文件
    target_files = []
    for client_data in document_list:
        for doc in client_data['documents']:
            if doc['file_type'] in ['pdf', 'powerpoint']:
                full_path = BASE_DIR / doc['filepath']
                if full_path.exists():
                    target_files.append({
                        'client_name': client_data['client_name'],
                        'filepath': full_path,
                        'relative_path': doc['filepath'],
                        'filename': doc['filename'],
                        'file_type': doc['file_type']
                    })
    
    print(f"\n找到 {len(target_files)} 个文档 (PDF: {sum(1 for f in target_files if f['file_type']=='pdf')}, PPTX: {sum(1 for f in target_files if f['file_type']=='powerpoint')})")
    
    # 提取数据
    results = []
    
    success_count = 0
    error_count = 0
    
    print(f"\n📊 开始提取数据...")
    for i, file_info in enumerate(target_files, 1):
        print(f"\n[{i:4d}/{len(target_files)}] 提取: {file_info['filename']}")
        print(f"          客户: {file_info['client_name']} | 类型: {file_info['file_type']}")
        
        if file_info['file_type'] == 'pdf':
            result = extract_pdf_text(file_info['filepath'])
        else:
            result = extract_pptx_text(file_info['filepath'])
        
        result['client_name'] = file_info['client_name']
        
        if result['success']:
            success_count += 1
            print(f"          ✅ 成功")
            if result.get('page_count') or result.get('slide_count'):
                count = result.get('page_count') or result.get('slide_count')
                print(f"          📄 页数: {count}")
            if result.get('summary'):
                print(f"          📝 {result['summary'][:50]}...")
        else:
            error_count += 1
            print(f"          ❌ 失败: {result.get('error', '未知错误')[:50]}")
        
        results.append(result)
    
    # 保存结果
    output_file = OUTPUT_DIR / 'pdf_pptx_extraction_results.json'
    
    print(f"\n\n💾 保存结果...")
    
    # 转换datetime类型
    def convert_to_serializable(obj):
        if isinstance(obj, dict):
            return {k: convert_to_serializable(v) for k, v in obj.items()}
        elif isinstance(obj, (list, tuple)):
            return [convert_to_serializable(item) for item in obj]
        elif isinstance(obj, datetime):
            return obj.isoformat()
        else:
            return obj
    
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(convert_to_serializable(results), f, ensure_ascii=False, indent=2)
    
    # 打印统计
    print("\n" + "=" * 60)
    print("✅ PDF/PPTX提取完成！")
    print("=" * 60)
    
    print(f"\n📊 统计:")
    print(f"   总文件数: {len(results)}")
    print(f"   成功: {success_count} ✅")
    print(f"   失败: {error_count} ❌")
    print(f"   成功率: {success_count/len(results)*100:.1f}%")
    
    print(f"\n💾 文件已保存:")
    print(f"   - 结果: {output_file}")

if __name__ == '__main__':
    main()
