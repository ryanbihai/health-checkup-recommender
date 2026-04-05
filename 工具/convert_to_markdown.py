#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件转换为Markdown格式脚本
将各种文件类型（Excel、PDF、PPT、Word、图片等）转换为md格式
保持原有文件夹结构和文件名
"""

import os
import sys
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Any, Optional

# 路径配置
BASE_DIR = Path(r'c:\IT\02 代理人营销工具\agent-customer-management')
INPUT_DIR = BASE_DIR / '待整理档案'
OUTPUT_DIR = BASE_DIR / '待整理档案md'

# 支持的文件类型
SUPPORTED_EXTENSIONS = {
    '.xlsx', '.xls',    # Excel
    '.pdf',             # PDF
    '.pptx',           # PowerPoint
    '.docx', '.doc',    # Word
    '.jpg', '.jpeg', '.png', '.gif', '.bmp',  # 图片
    '.xmind',          # XMind
    '.jsf',            # JSF
}


class MarkdownConverter:
    """Markdown转换器基类"""
    
    def __init__(self):
        self.success_count = 0
        self.error_count = 0
        self.skip_count = 0
    
    def convert(self, input_path: Path, output_path: Path) -> bool:
        """转换单个文件"""
        raise NotImplementedError
    
    def get_supported_extensions(self) -> set:
        """获取支持的文件扩展名"""
        raise NotImplementedError


class ExcelConverter(MarkdownConverter):
    """Excel转Markdown转换器"""
    
    def get_supported_extensions(self) -> set:
        return {'.xlsx', '.xls'}
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """转换Excel文件，每个worksheet转为一个md文件"""
        try:
            # 根据文件类型选择引擎
            if input_path.suffix.lower() == '.xlsx':
                # xlsx格式使用openpyxl
                return self._convert_xlsx(input_path, output_dir)
            else:
                # xls格式使用xlrd
                return self._convert_xls(input_path, output_dir)
                
        except Exception as e:
            print(f"      错误: {str(e)}")
            return False
    
    def _convert_xlsx(self, input_path: Path, output_dir: Path) -> bool:
        """转换xlsx文件"""
        import pandas as pd
        
        try:
            excel_file = pd.ExcelFile(str(input_path), engine='openpyxl')
            sheet_names = excel_file.sheet_names
            
            if len(sheet_names) == 1:
                output_path = output_dir / f"{input_path.stem}.md"
                self._convert_single_sheet_xlsx(input_path, sheet_names[0], output_path)
                return True
            else:
                for sheet_name in sheet_names:
                    safe_sheet_name = self._sanitize_filename(sheet_name)
                    if safe_sheet_name:
                        output_path = output_dir / f"{input_path.stem}_{safe_sheet_name}.md"
                    else:
                        output_path = output_dir / f"{input_path.stem}.md"
                    self._convert_single_sheet_xlsx(input_path, sheet_name, output_path)
                return True
                
        except Exception as e:
            print(f"      错误: {str(e)}")
            return False
    
    def _convert_xls(self, input_path: Path, output_dir: Path) -> bool:
        """转换xls文件（使用xlrd）"""
        try:
            import xlrd
            
            wb = xlrd.open_workbook(str(input_path))
            sheet_names = wb.sheet_names()
            
            if len(sheet_names) == 1:
                output_path = output_dir / f"{input_path.stem}.md"
                self._convert_single_sheet_xls(wb, wb.sheet_by_index(0), output_path)
                return True
            else:
                for sheet_name in sheet_names:
                    sheet = wb.sheet_by_name(sheet_name)
                    safe_sheet_name = self._sanitize_filename(sheet_name)
                    if safe_sheet_name:
                        output_path = output_dir / f"{input_path.stem}_{safe_sheet_name}.md"
                    else:
                        output_path = output_dir / f"{input_path.stem}.md"
                    self._convert_single_sheet_xls(wb, sheet, output_path)
                return True
                
        except ImportError:
            print(f"      ⚠️  未安装xlrd库，将尝试pandas处理")
            # 降级方案：尝试用pandas的xlrd引擎
            return self._convert_xls_with_pandas(input_path, output_dir)
        except Exception as e:
            print(f"      错误: {str(e)}")
            return False
    
    def _convert_xls_with_pandas(self, input_path: Path, output_dir: Path) -> bool:
        """使用pandas的xlrd引擎转换xls"""
        import pandas as pd
        
        try:
            # 尝试用pandas读取xls
            excel_file = pd.ExcelFile(str(input_path), engine='xlrd')
            sheet_names = excel_file.sheet_names
            
            for sheet_name in sheet_names:
                df = pd.read_excel(str(input_path), sheet_name=sheet_name, engine='xlrd', header=None)
                
                # 生成markdown
                md_content = []
                md_content.append(f"# {input_path.stem}\n")
                md_content.append(f"**工作表**: {sheet_name}\n")
                md_content.append(f"**源文件**: {input_path.name}\n")
                md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                md_content.append("\n---\n\n")
                
                if not df.empty:
                    md_content.append("## 数据表格\n\n")
                    
                    # 生成表格
                    header = '| ' + ' | '.join([str(i) for i in range(len(df.columns))]) + ' |'
                    separator = '|' + '|'.join([' --- ' for _ in df.columns]) + '|'
                    
                    md_content.append(header + '\n')
                    md_content.append(separator + '\n')
                    
                    for idx, row in df.iterrows():
                        values = []
                        for val in row:
                            if pd.isna(val):
                                values.append('')
                            else:
                                values.append(str(val).replace('\n', ' ').replace('|', '\\|'))
                        md_content.append('| ' + ' | '.join(values) + ' |\n')
                
                safe_sheet_name = self._sanitize_filename(sheet_name)
                if safe_sheet_name:
                    output_path = output_dir / f"{input_path.stem}_{safe_sheet_name}.md"
                else:
                    output_path = output_dir / f"{input_path.stem}.md"
                
                output_path.parent.mkdir(parents=True, exist_ok=True)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(''.join(md_content))
                
                self.success_count += 1
                print(f"      ✅ 已生成: {output_path.name}")
            
            return True
            
        except ImportError:
            print(f"      ❌ 需安装xlrd库: pip install xlrd")
            return False
        except Exception as e:
            print(f"      错误: {str(e)}")
            return False
    
    def _convert_single_sheet_xls(self, wb, sheet, output_path: Path):
        """转换单个xls工作表"""
        import xlrd
        from xlrd import xldate
        
        md_content = []
        md_content.append(f"# {output_path.stem.replace('_' + sheet.name, '')}\n")
        md_content.append(f"**工作表**: {sheet.name}\n")
        md_content.append(f"**行数**: {sheet.nrows}\n")
        md_content.append(f"**列数**: {sheet.ncols}\n")
        md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        md_content.append("\n---\n\n")
        
        if sheet.nrows > 0:
            md_content.append("## 数据表格\n\n")
            
            # 生成表头
            header_row = 0
            for idx in range(min(5, sheet.nrows)):
                row_values = [str(sheet.cell_value(idx, col)) for col in range(sheet.ncols)]
                row_str = ' '.join([v for v in row_values if v and v != 'nan'])
                if any(kw in row_str for kw in ['投保人', '被保险人', '险种', '保额', '保费', '生效']):
                    header_row = idx
                    break
            
            # 使用找到的标题行
            header_values = []
            for col in range(sheet.ncols):
                cell = sheet.cell(header_row, col)
                value = cell.value
                if isinstance(value, float) and cell.ctype == xlrd.XL_CELL_DATE:
                    # 日期格式
                    try:
                        date_tuple = xlrd.xldate_as_tuple(value, wb.datemode)
                        if date_tuple[0] > 0:
                            value = f"{date_tuple[0]}-{date_tuple[1]:02d}-{date_tuple[2]:02d}"
                        else:
                            value = str(value)
                    except:
                        value = str(value)
                else:
                    value = str(value) if value else ''
                header_values.append(value.replace('\n', ' ').strip())
            
            header = '| ' + ' | '.join(header_values) + ' |'
            separator = '|' + '|'.join([' --- ' for _ in header_values]) + '|'
            
            md_content.append(header + '\n')
            md_content.append(separator + '\n')
            
            # 添加数据行
            for row_idx in range(header_row + 1, sheet.nrows):
                values = []
                for col in range(sheet.ncols):
                    cell = sheet.cell(row_idx, col)
                    value = cell.value
                    
                    if cell.ctype == xlrd.XL_CELL_EMPTY or value == '':
                        values.append('')
                    elif cell.ctype == xlrd.XL_CELL_DATE:
                        # 日期格式
                        try:
                            date_tuple = xlrd.xldate_as_tuple(value, wb.datemode)
                            if date_tuple[0] > 0:
                                value = f"{date_tuple[0]}-{date_tuple[1]:02d}-{date_tuple[2]:02d}"
                            else:
                                value = str(value)
                        except:
                            value = str(value)
                    elif isinstance(value, float):
                        # 数字，保留两位小数
                        if value == int(value):
                            value = str(int(value))
                        else:
                            value = f"{value:.2f}"
                    else:
                        value = str(value).replace('\n', ' ').replace('|', '\\|')
                    
                    values.append(value)
                
                md_content.append('| ' + ' | '.join(values) + ' |\n')
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(''.join(md_content))
        
        self.success_count += 1
        print(f"      ✅ 已生成: {output_path.name}")
    
    def _convert_single_sheet_xlsx(self, input_path: Path, sheet_name: str, output_path: Path):
        """转换单个xlsx工作表为md"""
        import pandas as pd
        
        try:
            # 尝试找到标题行
            df_raw = pd.read_excel(str(input_path), sheet_name=sheet_name, header=None, engine='openpyxl')
            
            header_row = 0
            for idx in range(min(5, len(df_raw))):
                row = df_raw.iloc[idx]
                row_str = ' '.join([str(v) for v in row if pd.notna(v)])
                if any(kw in row_str for kw in ['投保人', '被保险人', '险种', '保额', '保费', '生效']):
                    header_row = idx
                    break
            
            # 重新读取，使用找到的标题行
            df = pd.read_excel(str(input_path), sheet_name=sheet_name, header=header_row, engine='openpyxl')
            
            # 生成markdown内容
            md_content = []
            md_content.append(f"# {input_path.stem}\n")
            md_content.append(f"**工作表**: {sheet_name}\n")
            md_content.append(f"**源文件**: {input_path.name}\n")
            md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            md_content.append("\n---\n\n")
            
            if not df.empty:
                # 添加表格
                md_content.append("## 数据表格\n\n")
                
                # 清理列名
                df.columns = [str(col).replace('\n', ' ').strip() if pd.notna(col) else '' for col in df.columns]
                
                # 生成表格头
                header = '| ' + ' | '.join(str(col) for col in df.columns) + ' |'
                separator = '|' + '|'.join([' --- ' for _ in df.columns]) + '|'
                
                md_content.append(header + '\n')
                md_content.append(separator + '\n')
                
                # 添加数据行
                for idx, row in df.iterrows():
                    values = []
                    for val in row:
                        if pd.isna(val):
                            values.append('')
                        elif isinstance(val, (int, float)):
                            values.append(str(val))
                        else:
                            values.append(str(val).replace('\n', ' ').replace('|', '\\|'))
                    md_content.append('| ' + ' | '.join(values) + ' |\n')
            
            # 保存文件
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(''.join(md_content))
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
    
    def _sanitize_filename(self, name: str) -> str:
        """清理文件名中的非法字符"""
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            name = name.replace(char, '_')
        return name.strip()


class PDFConverter(MarkdownConverter):
    """PDF转Markdown转换器"""
    
    def get_supported_extensions(self) -> set:
        return {'.pdf'}
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """转换PDF文件为md"""
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            # 尝试使用PyMuPDF
            try:
                import fitz
                doc = fitz.open(str(input_path))
                
                md_content = []
                md_content.append(f"# {input_path.stem}\n")
                md_content.append(f"**源文件**: {input_path.name}\n")
                md_content.append(f"**页数**: {len(doc)}\n")
                md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                md_content.append("\n---\n\n")
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text = page.get_text()
                    if text.strip():
                        md_content.append(f"## 第 {page_num + 1} 页\n\n")
                        md_content.append(text)
                        md_content.append("\n\n")
                
                doc.close()
                
            except ImportError:
                # 尝试使用pdfplumber
                import pdfplumber
                
                md_content = []
                md_content.append(f"# {input_path.stem}\n")
                md_content.append(f"**源文件**: {input_path.name}\n")
                md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                md_content.append("\n---\n\n")
                
                with pdfplumber.open(str(input_path)) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text:
                            md_content.append(f"## 第 {page_num} 页\n\n")
                            md_content.append(text)
                            md_content.append("\n\n")
            
            # 保存文件
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(''.join(md_content))
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            return True
            
        except ImportError:
            print(f"      ⚠️  未安装PDF处理库，将记录文件信息")
            self._create_placeholder(input_path, output_path, "PDF文档")
            return True
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False
    
    def _create_placeholder(self, input_path: Path, output_path: Path, doc_type: str):
        """创建占位符md文件"""
        md_content = f"""# {input_path.stem}

**源文件**: {input_path.name}
**文件类型**: {doc_type}
**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

---

> ⚠️ 无法提取内容，仅记录文件信息
> 
> 请安装必要的库来处理此文件格式:
> - PDF: pip install PyMuPDF 或 pip install pdfplumber
"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        self.success_count += 1


class PPTXConverter(MarkdownConverter):
    """PowerPoint转Markdown转换器"""
    
    def get_supported_extensions(self) -> set:
        return {'.pptx'}
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """转换PPTX文件为md"""
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            from pptx import Presentation
            
            prs = Presentation(str(input_path))
            
            md_content = []
            md_content.append(f"# {input_path.stem}\n")
            md_content.append(f"**源文件**: {input_path.name}\n")
            md_content.append(f"**幻灯片数**: {len(prs.slides)}\n")
            md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            md_content.append("\n---\n\n")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                
                if slide_texts:
                    md_content.append(f"## 幻灯片 {slide_num}\n\n")
                    md_content.append('\n\n'.join(slide_texts))
                    md_content.append("\n\n")
            
            # 保存文件
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(''.join(md_content))
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            return True
            
        except ImportError:
            print(f"      ⚠️  未安装PPTX处理库")
            self._create_placeholder(input_path, output_path)
            return True
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False
    
    def _create_placeholder(self, input_path: Path, output_path: Path):
        """创建占位符md文件"""
        md_content = f"""# {input_path.stem}

**源文件**: {input_path.name}
**文件类型**: PowerPoint演示文稿
**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

---

> ⚠️ 无法提取内容，仅记录文件信息
> 
> 请安装必要的库来处理此文件格式:
> - PPTX: pip install python-pptx
"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        self.success_count += 1


class DOCXConverter(MarkdownConverter):
    """Word转Markdown转换器"""
    
    def get_supported_extensions(self) -> set:
        return {'.docx', '.doc'}
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """转换Word文档为md"""
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            if input_path.suffix.lower() == '.docx':
                from docx import Document
                
                doc = Document(str(input_path))
                
                md_content = []
                md_content.append(f"# {input_path.stem}\n")
                md_content.append(f"**源文件**: {input_path.name}\n")
                md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                md_content.append("\n---\n\n")
                
                for para in doc.paragraphs:
                    if para.text.strip():
                        md_content.append(para.text)
                        md_content.append("\n\n")
                
                # 保存文件
                output_path.parent.mkdir(parents=True, exist_ok=True)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(''.join(md_content))
                
                self.success_count += 1
                print(f"      ✅ 已生成: {output_path.name}")
                return True
            else:
                # 处理.doc文件（RTF格式）
                return self._convert_doc_file(input_path, output_dir)
                
        except ImportError:
            print(f"      ⚠️  未安装Word处理库")
            self._create_placeholder(input_path, output_path)
            return True
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False
    
    def _convert_doc_file(self, input_path: Path, output_dir: Path) -> bool:
        """处理.doc文件"""
        import re
        
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # 移除RTF格式标签
            text = re.sub(r'\\\{|\}|<[^>]+>', ' ', content)
            text = re.sub(r'\\[a-z]+\d+\s', ' ', text)
            text = re.sub(r'\s+', ' ', text)
            
            md_content = []
            md_content.append(f"# {input_path.stem}\n")
            md_content.append(f"**源文件**: {input_path.name}\n")
            md_content.append(f"**文件类型**: Word文档 (.doc)\n")
            md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            md_content.append("\n---\n\n")
            md_content.append(text.strip())
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(''.join(md_content))
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            return True
            
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False
    
    def _create_placeholder(self, input_path: Path, output_path: Path):
        """创建占位符md文件"""
        md_content = f"""# {input_path.stem}

**源文件**: {input_path.name}
**文件类型**: Word文档
**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

---

> ⚠️ 无法提取内容，仅记录文件信息
> 
> 请安装必要的库来处理此文件格式:
> - DOCX: pip install python-docx
"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        self.success_count += 1


class ImageConverter(MarkdownConverter):
    """图片转Markdown转换器"""
    
    def get_supported_extensions(self) -> set:
        return {'.jpg', '.jpeg', '.png', '.gif', '.bmp'}
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """转换图片为md（记录图片信息）"""
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            # 获取文件大小
            file_size = input_path.stat().st_size
            if file_size > 1024 * 1024:
                size_str = f"{file_size / (1024*1024):.2f} MB"
            elif file_size > 1024:
                size_str = f"{file_size / 1024:.2f} KB"
            else:
                size_str = f"{file_size} bytes"
            
            # 获取图片尺寸（如果可能）
            dimensions = ""
            try:
                from PIL import Image
                with Image.open(input_path) as img:
                    dimensions = f"{img.width} x {img.height}"
            except:
                pass
            
            md_content = []
            md_content.append(f"# {input_path.stem}\n")
            md_content.append(f"**源文件**: {input_path.name}\n")
            md_content.append(f"**文件类型**: 图片\n")
            md_content.append(f"**文件大小**: {size_str}\n")
            if dimensions:
                md_content.append(f"**图片尺寸**: {dimensions}\n")
            md_content.append(f"**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            md_content.append("\n---\n\n")
            md_content.append(f"![{input_path.name}]({input_path.absolute()})\n")
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(''.join(md_content))
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            return True
            
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False


class XMindConverter(MarkdownConverter):
    """XMind转Markdown转换器"""
    
    def get_supported_extensions(self) -> set:
        return {'.xmind'}
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """转换XMind文件为md（记录文件信息）"""
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            file_size = input_path.stat().st_size
            if file_size > 1024 * 1024:
                size_str = f"{file_size / (1024*1024):.2f} MB"
            elif file_size > 1024:
                size_str = f"{file_size / 1024:.2f} KB"
            else:
                size_str = f"{file_size} bytes"
            
            md_content = f"""# {input_path.stem}

**源文件**: {input_path.name}
**文件类型**: XMind思维导图
**文件大小**: {size_str}
**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

---

> ⚠️ XMind文件需要使用XMind软件打开查看
> 
> 原始文件路径: `{input_path.absolute()}`

## 文件说明

此文件为XMind格式的思维导图文件，包含以下可能的结构：
- 中心主题
- 分支主题
- 自由主题
- 标签
- 备注
- 资源链接

如需查看或编辑，请使用XMind软件打开源文件。
"""
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            return True
            
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False


class JSFConverter(MarkdownConverter):
    """JSF文件转Markdown转换器"""
    
    def get_supported_extensions(self) -> set:
        return {'.jsf'}
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """转换JSF文件为md（记录文件信息）"""
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            file_size = input_path.stat().st_size
            if file_size > 1024 * 1024:
                size_str = f"{file_size / (1024*1024):.2f} MB"
            elif file_size > 1024:
                size_str = f"{file_size / 1024:.2f} KB"
            else:
                size_str = f"{file_size} bytes"
            
            md_content = f"""# {input_path.stem}

**源文件**: {input_path.name}
**文件类型**: JSF文件
**文件大小**: {size_str}
**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

---

> ⚠️ 无法识别JSF文件格式，仅记录文件信息
> 
> 原始文件路径: `{input_path.absolute()}`
"""
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            return True
            
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False


class UnknownConverter(MarkdownConverter):
    """未知文件类型转换器"""
    
    def get_supported_extensions(self) -> set:
        return set()
    
    def convert(self, input_path: Path, output_dir: Path) -> bool:
        """记录未知文件类型"""
        output_path = output_dir / f"{input_path.stem}.md"
        
        try:
            file_size = input_path.stat().st_size
            if file_size > 1024 * 1024:
                size_str = f"{file_size / (1024*1024):.2f} MB"
            elif file_size > 1024:
                size_str = f"{file_size / 1024:.2f} KB"
            else:
                size_str = f"{file_size} bytes"
            
            md_content = f"""# {input_path.stem}

**源文件**: {input_path.name}
**文件类型**: {input_path.suffix.upper()} 文件
**文件大小**: {size_str}
**转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

---

> ⚠️ 不支持的文件类型，仅记录文件信息
> 
> 原始文件路径: `{input_path.absolute()}`
"""
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            self.success_count += 1
            print(f"      ✅ 已生成: {output_path.name}")
            return True
            
        except Exception as e:
            print(f"      ❌ 错误: {str(e)}")
            self.error_count += 1
            return False


def get_converter(file_extension: str) -> MarkdownConverter:
    """根据文件扩展名获取对应的转换器"""
    converters = {
        '.xlsx': ExcelConverter(),
        '.xls': ExcelConverter(),
        '.pdf': PDFConverter(),
        '.pptx': PPTXConverter(),
        '.docx': DOCXConverter(),
        '.doc': DOCXConverter(),
        '.jpg': ImageConverter(),
        '.jpeg': ImageConverter(),
        '.png': ImageConverter(),
        '.gif': ImageConverter(),
        '.bmp': ImageConverter(),
        '.xmind': XMindConverter(),
        '.jsf': JSFConverter(),
    }
    return converters.get(file_extension.lower(), UnknownConverter())


def scan_files(input_dir: Path) -> List[Path]:
    """扫描目录下所有支持的文件"""
    files = []
    for ext in SUPPORTED_EXTENSIONS:
        files.extend(input_dir.rglob(f'*{ext}'))
        files.extend(input_dir.rglob(f'*{ext.upper()}'))
    return sorted(files)


def main():
    """主函数"""
    print("=" * 70)
    print("文件转换为Markdown工具")
    print("=" * 70)
    
    print(f"\n📂 源目录: {INPUT_DIR}")
    print(f"📁 输出目录: {OUTPUT_DIR}")
    
    # 创建输出目录
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    # 扫描文件
    print(f"\n🔍 扫描文件中...")
    files = scan_files(INPUT_DIR)
    
    if not files:
        print("❌ 未找到支持的文件")
        return
    
    print(f"✅ 找到 {len(files)} 个文件\n")
    
    # 按文件夹分组统计
    folder_stats = {}
    for f in files:
        folder = f.parent.relative_to(INPUT_DIR)
        folder_str = str(folder) if str(folder) != '.' else '(根目录)'
        if folder_str not in folder_stats:
            folder_stats[folder_str] = {'count': 0, 'files': []}
        folder_stats[folder_str]['count'] += 1
        folder_stats[folder_str]['files'].append(f)
    
    print("📊 文件分布:")
    for folder, stats in sorted(folder_stats.items())[:10]:
        print(f"   {folder}/ ({stats['count']} 个文件)")
    if len(folder_stats) > 10:
        print(f"   ... 还有 {len(folder_stats) - 10} 个文件夹")
    print()
    
    # 统计文件类型
    ext_stats = {}
    for f in files:
        ext = f.suffix.lower()
        ext_stats[ext] = ext_stats.get(ext, 0) + 1
    
    print("📈 文件类型统计:")
    for ext, count in sorted(ext_stats.items(), key=lambda x: -x[1]):
        print(f"   {ext}: {count} 个")
    print()
    
    # 开始转换
    print("=" * 70)
    print("开始转换...")
    print("=" * 70)
    
    total_success = 0
    total_error = 0
    total_skip = 0
    
    for idx, file_path in enumerate(files, 1):
        print(f"\n[{idx:4d}/{len(files)}] {file_path.name}")
        print(f"         路径: {file_path.parent.relative_to(INPUT_DIR)}")
        
        # 计算输出目录（保持原有文件夹结构）
        relative_path = file_path.parent.relative_to(INPUT_DIR)
        output_subdir = OUTPUT_DIR / relative_path
        
        # 获取转换器
        converter = get_converter(file_path.suffix)
        
        # 执行转换
        try:
            success = converter.convert(file_path, output_subdir)
            if success:
                total_success += 1
            else:
                total_error += 1
        except Exception as e:
            print(f"      ❌ 转换异常: {str(e)}")
            total_error += 1
    
    # 打印统计
    print("\n" + "=" * 70)
    print("转换完成！")
    print("=" * 70)
    print(f"\n📊 统计:")
    print(f"   总文件数: {len(files)}")
    print(f"   成功: {total_success} ✅")
    print(f"   失败: {total_error} ❌")
    print(f"   成功率: {total_success/len(files)*100:.1f}%")
    print(f"\n📁 输出目录: {OUTPUT_DIR}")
    
    # 统计输出文件
    md_files = list(OUTPUT_DIR.rglob('*.md'))
    print(f"   生成MD文件: {len(md_files)} 个")
    
    print("\n✨ 转换完成！")


if __name__ == '__main__':
    main()
